from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from .common import clean_text, extraction_failure


class PptxExtractor:
    name = "python-pptx"
    version = "1"

    def extract(self, path: Path, max_pages: int):
        try:
            presentation = Presentation(str(path))
            slides = presentation.slides if max_pages <= 0 else presentation.slides[:max_pages]
            parts: list[str] = []
            slide_count = 0
            for slide_number, slide in enumerate(slides, 1):
                slide_count += 1
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        value = clean_text(shape.text)
                        if value:
                            slide_parts.append(value)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            values = [clean_text(cell.text) for cell in row.cells]
                            slide_parts.append(" | ".join(value for value in values if value))
                try:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        slide_parts.append(clean_text(notes))
                except Exception:
                    pass
                if slide_parts:
                    parts.append(f"Slide {slide_number}: " + " ".join(slide_parts))
            text = clean_text("\n".join(parts))
            metadata = {"slide_count": len(presentation.slides), "slides_extracted": slide_count}
            if not text:
                return "", metadata, extraction_failure("NO_EXTRACTABLE_TEXT", "The PPTX file contains no extractable text.")
            return text, metadata, None
        except Exception as exc:
            return "", {}, extraction_failure("PPTX_EXTRACTION_FAILED", str(exc))
