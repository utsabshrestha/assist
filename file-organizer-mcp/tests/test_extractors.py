import json
from docx import Document
from pptx import Presentation
from file_organizer_mcp.extractors import build_extractor_registry


def test_registry_has_phase_two_extensions():
    assert set(build_extractor_registry()) == {".pdf", ".txt", ".md", ".docx", ".json", ".html", ".htm", ".xml", ".pptx", ".epub"}


def test_plain_text_and_markdown(tmp_path):
    registry = build_extractor_registry()
    for suffix in (".txt", ".md"):
        path = tmp_path / f"sample{suffix}"
        path.write_text("Alpha   beta\nGamma", encoding="utf-8")
        text, _, error = registry[suffix].extract(path, 0)
        assert error is None
        assert text == "Alpha beta Gamma"


def test_json_preserves_keys_and_values(tmp_path):
    path = tmp_path / "sample.json"
    path.write_text(json.dumps({"department": "Accessibility", "count": 3}), encoding="utf-8")
    text, metadata, error = build_extractor_registry()[".json"].extract(path, 0)
    assert error is None
    assert "department" in text and "Accessibility" in text
    assert metadata["root_type"] == "dict"


def test_html_excludes_script_and_style(tmp_path):
    path = tmp_path / "sample.html"
    path.write_text("<html><head><style>hidden</style></head><body><h1>Visible</h1><script>secret</script></body></html>", encoding="utf-8")
    text, _, error = build_extractor_registry()[".html"].extract(path, 0)
    assert error is None
    assert "Visible" in text and "secret" not in text and "hidden" not in text


def test_xml_includes_tags_and_values(tmp_path):
    path = tmp_path / "sample.xml"
    path.write_text("<report><title>Annual Review</title></report>", encoding="utf-8")
    text, _, error = build_extractor_registry()[".xml"].extract(path, 0)
    assert error is None
    assert "report" in text and "Annual Review" in text


def test_docx_paragraphs_and_tables(tmp_path):
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Project overview")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Utsab"
    document.save(path)
    text, metadata, error = build_extractor_registry()[".docx"].extract(path, 0)
    assert error is None
    assert "Project overview" in text and "Owner" in text and "Utsab" in text
    assert metadata["table_count"] == 1


def test_pptx_text(tmp_path):
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    slide.placeholders[1].text = "Accessibility progress"
    presentation.save(path)
    text, metadata, error = build_extractor_registry()[".pptx"].extract(path, 0)
    assert error is None
    assert "Quarterly Review" in text and "Accessibility progress" in text
    assert metadata["slide_count"] == 1
